from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # white
C_HEADER   = RGBColor(0x1A, 0x53, 0x76)   # dark blue
C_ACCENT   = RGBColor(0xE8, 0x4C, 0x3B)   # red/Italian
C_GREEN    = RGBColor(0x2E, 0x7D, 0x4F)   # green/Italian
C_BODY     = RGBColor(0x33, 0x33, 0x33)   # dark grey
C_LIGHT    = RGBColor(0xF0, 0xF4, 0xF8)   # light blue bg
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_YELLOW   = RGBColor(0xFF, 0xF0, 0x80)

W = Inches(10)
H = Inches(5.625)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # truly blank


def add_rect(slide, x, y, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_BODY
    return txb


def slide_title_banner(slide, title, subtitle=None):
    # full bg
    add_rect(slide, 0, 0, W, H, C_HEADER)
    # Italian flag stripe
    stripe_h = Inches(0.12)
    add_rect(slide, 0, 0, W // 3, stripe_h, C_GREEN)
    add_rect(slide, W // 3, 0, W // 3, stripe_h, C_WHITE)
    add_rect(slide, W * 2 // 3, 0, W // 3, stripe_h, C_ACCENT)
    # lesson number chip
    chip_w, chip_h = Inches(1.4), Inches(0.55)
    chip_x = Inches(0.4)
    chip_y = Inches(1.2)
    add_rect(slide, chip_x, chip_y, chip_w, chip_h, C_ACCENT)
    add_text(slide, "Lezione 19", chip_x, chip_y + Inches(0.04),
             chip_w, chip_h, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # main title
    add_text(slide, title, Inches(0.4), Inches(2.0),
             Inches(9.2), Inches(1.4), size=46, bold=True,
             color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(3.5),
                 Inches(9.2), Inches(0.8), size=22,
                 color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.LEFT)
    # bottom tag
    add_text(slide, "Italiano per stranieri", Inches(0.4), Inches(4.9),
             Inches(9.2), Inches(0.5), size=14,
             color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.LEFT)


def slide_content(slide, title, badge=None, badge_color=None, items=None,
                  note=None, two_col=None):
    """Generic content slide."""
    add_rect(slide, 0, 0, W, H, C_BG)
    # top bar
    add_rect(slide, 0, 0, W, Inches(0.9), C_HEADER)
    add_text(slide, title, Inches(0.3), Inches(0.1),
             Inches(7.5), Inches(0.75), size=24, bold=True, color=C_WHITE)
    if badge:
        bc = badge_color or C_ACCENT
        bw = Inches(1.6)
        add_rect(slide, W - bw - Inches(0.2), Inches(0.18), bw, Inches(0.54), bc)
        add_text(slide, badge, W - bw - Inches(0.2), Inches(0.2),
                 bw, Inches(0.5), size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # body items
    if items and not two_col:
        y = Inches(1.1)
        for item in items:
            bullet, text, sz, clr, ind = item
            bx = Inches(0.4) + Inches(ind * 0.3)
            bw2 = Inches(9.2) - Inches(ind * 0.3)
            if bullet:
                add_text(slide, bullet, bx, y, Inches(0.35), Inches(0.42),
                         size=sz, bold=True, color=clr or C_ACCENT)
                add_text(slide, text, bx + Inches(0.4), y, bw2 - Inches(0.4),
                         Inches(0.42), size=sz, color=clr or C_BODY)
            else:
                add_text(slide, text, bx, y, bw2, Inches(0.42),
                         size=sz, bold=False, color=clr or C_BODY)
            y += Inches(0.44)
    if two_col:
        # two_col = (left_items, right_items, left_header, right_header)
        left, right, lh, rh = two_col
        mid = W // 2
        col_w = Inches(4.4)
        # left header
        add_rect(slide, Inches(0.3), Inches(1.05), col_w, Inches(0.5),
                 C_ACCENT)
        add_text(slide, lh, Inches(0.3), Inches(1.08), col_w, Inches(0.44),
                 size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # right header
        add_rect(slide, mid + Inches(0.3), Inches(1.05), col_w, Inches(0.5),
                 C_GREEN)
        add_text(slide, rh, mid + Inches(0.3), Inches(1.08), col_w, Inches(0.44),
                 size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        y_l = y_r = Inches(1.7)
        for txt in left:
            add_text(slide, txt, Inches(0.5), y_l, col_w - Inches(0.2),
                     Inches(0.44), size=17, color=C_BODY)
            y_l += Inches(0.46)
        for txt in right:
            add_text(slide, txt, mid + Inches(0.5), y_r, col_w - Inches(0.2),
                     Inches(0.44), size=17, color=C_BODY)
            y_r += Inches(0.46)
    if note:
        note_y = Inches(5.0)
        add_rect(slide, Inches(0.3), note_y, Inches(9.4), Inches(0.45), C_YELLOW)
        add_text(slide, note, Inches(0.5), note_y + Inches(0.04),
                 Inches(9.0), Inches(0.38), size=13, color=C_BODY)


# ─── Slide 1: Title ───────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_title_banner(sl,
    "Passato Prossimo",
    "Come parlare del passato in italiano")

# ─── Slide 2: Che cos'è il Passato Prossimo ───────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "Che cos'è il Passato Prossimo?",
    badge="CORE", badge_color=C_ACCENT,
    items=[
        ("", "Il passato prossimo descrive azioni avvenute nel passato recente.", 18, None, 0),
        ("", "", 10, None, 0),
        ("📌", "Formula:", 18, C_HEADER, 0),
        ("", "", 10, None, 0),
        ("", "AUSILIARE  +  PARTICIPIO PASSATO", 26, C_ACCENT, 1),
        ("", "", 10, None, 0),
        ("✔", "avere  (ho, hai, ha, abbiamo, avete, hanno)", 17, C_BODY, 1),
        ("✔", "essere  (sono, sei, è, siamo, siete, sono)", 17, C_BODY, 1),
    ])

# ─── Slide 3: ho mangiato ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "Ausiliare AVERE — ho mangiato",
    badge="CORE", badge_color=C_ACCENT,
    items=[
        ("io",    "ho mangiato",         19, None, 0),
        ("tu",    "hai mangiato",        19, None, 0),
        ("lui/lei","ha mangiato",        19, None, 0),
        ("noi",   "abbiamo mangiato",    19, None, 0),
        ("voi",   "avete mangiato",      19, None, 0),
        ("loro",  "hanno mangiato",      19, None, 0),
    ],
    note='💬 Esempio: "Ho mangiato la pizza ieri sera."')

# ─── Slide 4: sono andato ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "Ausiliare ESSERE — sono andato/a",
    badge="CORE", badge_color=C_ACCENT,
    items=[
        ("io",     "sono andato / andata",       19, None, 0),
        ("tu",     "sei andato / andata",         19, None, 0),
        ("lui/lei","è andato / andata",           19, None, 0),
        ("noi",    "siamo andati / andate",       19, None, 0),
        ("voi",    "siete andati / andate",       19, None, 0),
        ("loro",   "sono andati / andate",        19, None, 0),
    ],
    note='⚠️ Con ESSERE il participio concorda con il soggetto  (andato / andata / andati / andate)')

# ─── Slide 5: AVERE vs ESSERE ─────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "AVERE o ESSERE?",
    badge="CORE", badge_color=C_ACCENT,
    two_col=(
        ["ho mangiato", "ho bevuto", "ho comprato", "ho studiato", "ho lavorato"],
        ["sono andato", "sono venuto", "sono rimasto", "sono partito", "mi sono alzato"],
        "AVERE  →  verbi transitivi",
        "ESSERE  →  moto/stato/riflessivi",
    ))

# ─── Slide 6: ieri ────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "Espressioni di tempo: ieri",
    badge="EXTRA", badge_color=C_GREEN,
    items=[
        ("🕐", "ieri  =  yesterday", 22, C_HEADER, 0),
        ("", "", 8, None, 0),
        ("▸", "Ieri ho studiato molto.", 18, None, 0),
        ("▸", "Ieri sera sono andato al cinema.", 18, None, 0),
        ("▸", "Ieri mattina ho fatto colazione con gli amici.", 18, None, 0),
        ("", "", 8, None, 0),
        ("📎", "ieri mattina  /  ieri pomeriggio  /  ieri sera", 15, C_BODY, 1),
    ])

# ─── Slide 7: settimana scorsa ───────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "Espressioni di tempo: settimana scorsa",
    badge="EXTRA", badge_color=C_GREEN,
    items=[
        ("🗓️", "settimana scorsa  =  last week", 22, C_HEADER, 0),
        ("", "", 8, None, 0),
        ("▸", "La settimana scorsa sono andato a Roma.", 18, None, 0),
        ("▸", "Settimana scorsa ho lavorato tanto.", 18, None, 0),
        ("", "", 8, None, 0),
        ("📎", "Simili:", 15, C_HEADER, 0),
        ("•", "mese scorso  /  anno scorso  /  lunedì scorso", 16, C_BODY, 1),
    ])

# ─── Slide 8: già / ancora ───────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "già  /  non... ancora",
    badge="EXTRA", badge_color=C_GREEN,
    items=[
        ("✅", "già  =  already  →  azione completata", 18, C_GREEN, 0),
        ("▸", "Ho già mangiato.  (I already ate.)", 17, None, 1),
        ("▸", "Sei già andato a Firenze?", 17, None, 1),
        ("", "", 8, None, 0),
        ("❌", "non... ancora  =  not yet  →  azione non completata", 18, C_ACCENT, 0),
        ("▸", "Non ho ancora mangiato.  (I haven't eaten yet.)", 17, None, 1),
        ("▸", "Non sono ancora andato a Firenze.", 17, None, 1),
    ])

# ─── Slide 9: Esercizio ───────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_content(sl,
    "Esercizio pratico",
    items=[
        ("", "Completa le frasi con il passato prossimo:", 18, C_HEADER, 0),
        ("", "", 6, None, 0),
        ("1.", "Ieri io _______ (mangiare) la pasta.", 18, None, 0),
        ("2.", "La settimana scorsa noi _______ (andare) al mare.", 18, None, 0),
        ("3.", "Ho _______ finito i compiti.  (già / ancora)", 18, None, 0),
        ("4.", "Non sono _______ stato a Venezia.  (già / ancora)", 18, None, 0),
    ],
    note="Risposte: 1. ho mangiato  2. siamo andati  3. già  4. ancora")

# ─── Slide 10: Riepilogo ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, C_HEADER)
add_rect(sl, 0, 0, W, Inches(0.12), C_GREEN)
add_text(sl, "Riepilogo", Inches(0.4), Inches(0.25),
         Inches(9.2), Inches(0.75), size=30, bold=True, color=C_WHITE)
add_rect(sl, Inches(0.4), Inches(1.1), Inches(9.2), Inches(3.6),
         RGBColor(0x0D, 0x3A, 0x5C))
checks = [
    "✅  Passato prossimo  =  ausiliare  +  participio passato",
    "✅  AVERE  →  verbi transitivi  (ho mangiato, ho bevuto…)",
    "✅  ESSERE  →  moto/stato/riflessivi  (sono andato, sono venuto…)",
    "✅  Con ESSERE il participio concorda col soggetto",
    "✅  Parole chiave: ieri, settimana scorsa, già, non...ancora",
]
for i, c in enumerate(checks):
    add_text(sl, c, Inches(0.7), Inches(1.2) + Inches(i * 0.68),
             Inches(8.8), Inches(0.6), size=17, color=C_WHITE)

out = "/home/user/LadyRadioApp/lezione19_passato_prossimo.pptx"
prs.save(out)
print("Saved:", out)
